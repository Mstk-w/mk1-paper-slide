import json
import re
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Configuration (A3 Size) ---
SLIDE_WIDTH_CM = 42.0
SLIDE_HEIGHT_CM = 29.7
MARGIN_CM = 1.5
HEADER_HEIGHT_CM = 2.5
BOX_GAP_CM = 0.8
COL_GAP_CM = 1.0

# --- Color Palette (Modern/Premium) ---
# Darker Navy for professionalism
COLOR_MAIN = RGBColor(0, 51, 102)      
COLOR_ACCENT = RGBColor(218, 165, 32)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_MAIN = RGBColor(33, 33, 33)
COLOR_TEXT_MUTED = RGBColor(100, 100, 100)
COLOR_BORDER = RGBColor(220, 220, 220) # Light Gray

# Fonts
FONT_NAME_BODY = "Meiryo UI"
FONT_NAME_BOLD = "Meiryo UI" 

def create_a3_slide(json_data, output_filename="output_slide.pptx"):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_WIDTH_CM)
    prs.slide_height = Cm(SLIDE_HEIGHT_CM)
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_WHITE

    # --- Header ---
    _draw_header(slide, json_data)

    # --- Layout Calculations ---
    # Adjust top margin
    content_top = MARGIN_CM + HEADER_HEIGHT_CM + 0.5
    content_height = SLIDE_HEIGHT_CM - content_top - MARGIN_CM
    
    total_content_width = SLIDE_WIDTH_CM - 2*MARGIN_CM
    col_width = (total_content_width - COL_GAP_CM) / 2
    
    left_x = MARGIN_CM
    right_x = MARGIN_CM + col_width + COL_GAP_CM

    # --- Content Processing ---
    content = json_data.get("content", [])
    
    left_items = []
    right_items = []

    if isinstance(content, dict):
        mapping = {
            "box1": "left", "box2": "left", "box3": "left", "box4": "left",
            "box5": "right", "box6": "right", "box7": "right", "box8": "right"
        }
        for k, v in content.items():
            side = "left"
            if "box5" in k or "box6" in k or "box7" in k or "box8" in k: side = "right"
            label = "Section"
            if "background" in k: label = "背景"
            elif "necessity" in k: label = "課題"
            elif "plan" in k: label = "施策"
            
            item = {"label": label, "text": v}
            if side == "left": left_items.append(item)
            else: right_items.append(item)
    else:
        for item in content:
            if item.get("column") == "left":
                left_items.append(item)
            elif item.get("column") == "right":
                right_items.append(item)

    # --- Draw Columns ---
    _draw_dynamic_column(slide, left_items, left_x, content_top, col_width, content_height)
    _draw_dynamic_column(slide, right_items, right_x, content_top, col_width, content_height)

    prs.save(output_filename)
    print(f"Generated: {output_filename}")


def _draw_header(slide, data):
    # Title Fitting Logic
    title_text = data.get('theme', 'Untitled')
    
    # Estimate size: if long, reduce font
    font_size = 32
    if len(title_text) > 20: font_size = 28
    if len(title_text) > 30: font_size = 24
    if len(title_text) > 40: font_size = 20
    
    title_box = slide.shapes.add_textbox(
        Cm(MARGIN_CM), Cm(MARGIN_CM), 
        Cm(SLIDE_WIDTH_CM - 2*MARGIN_CM), Cm(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.name = FONT_NAME_BOLD
    p.font.color.rgb = COLOR_MAIN
    
    # Subtitle / Department
    sub_box = slide.shapes.add_textbox(
        Cm(MARGIN_CM), Cm(MARGIN_CM + 1.5), 
        Cm(SLIDE_WIDTH_CM - 2*MARGIN_CM), Cm(1.0)
    )
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = data.get('department', '')
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.font.name = FONT_NAME_BODY

    # Accent Line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Cm(MARGIN_CM), Cm(MARGIN_CM + HEADER_HEIGHT_CM), 
        Cm(SLIDE_WIDTH_CM - 2*MARGIN_CM), Cm(0.05) # Thin line
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


def _draw_dynamic_column(slide, items, x, y, w, total_h):
    count = len(items)
    if count == 0: return

    total_gap = (count - 1) * BOX_GAP_CM
    avail_h = total_h - total_gap
    item_h = avail_h / count
    
    current_y = y
    
    for item in items:
        layout = item.get("layout_type", "text")
        if layout == "flow_horizontal":
            _draw_flow_horizontal(slide, x, current_y, w, item_h, item.get("label", ""), item.get("text", ""))
        else:
            _draw_section(slide, x, current_y, w, item_h, item.get("label", ""), item.get("text", ""))
        current_y += item_h + BOX_GAP_CM


def _draw_section(slide, x, y, w, h, label, text):
    # Border Box (Light Gray)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(x), Cm(y),
        Cm(w), Cm(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE # White bg
    box.line.color.rgb = COLOR_BORDER
    box.line.width = Pt(1.0) # Thin border

    header_h = 1.0
    
    # Accent Bar (Inside box)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(x + 0.1), Cm(y + 0.15),
        Cm(0.15), Cm(header_h - 0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_MAIN
    bar.line.fill.background()
    
    # Label Text (Heading: 20pt)
    label_box = slide.shapes.add_textbox(
        Cm(x + 0.4), Cm(y),
        Cm(w - 0.4), Cm(header_h)
    )
    p = label_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(20) # Requested 20pt
    p.font.bold = True
    p.font.name = FONT_NAME_BOLD
    p.font.color.rgb = COLOR_MAIN
    p.alignment = PP_ALIGN.LEFT
    label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 2. Content Body
    content_y = y + header_h
    content_h = h - header_h - 0.2
    
    # Auto-scaling font logic
    # Base 14pt, Scale to 12pt if needed
    char_count = len(text)
    font_size = 14
    if char_count > 250: font_size = 12
    
    text_box = slide.shapes.add_textbox(
        Cm(x + 0.4), Cm(content_y),
        Cm(w - 0.6), Cm(content_h)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    
    # Parse and add text
    _add_formatted_text(tf, text, font_size)


def _add_formatted_text(text_frame, raw_text, font_size_pt):
    lines = raw_text.split('\n')
    for line in lines:
        line = line.strip()
        if not line: continue
        
        p = text_frame.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        p.line_spacing = 1.2 # Requested 1.2
        
        clean_text = line
        # Use full-width bullet for aesthetics if desired, but sticking to text consistency
        if line.startswith("・") or line.startswith("-") or line.startswith("●"):
            clean_text = line[1:].strip()
            clean_text = "・" + clean_text # Enforce Japanese bullet
            
        parts = re.split(r'(\*\*.*?\*\*)', clean_text)
        
        for part in parts:
            if not part: continue
            
            is_bold = False
            content = part
            
            if part.startswith("**") and part.endswith("**"):
                is_bold = True
                content = part[2:-2]
                
            run = p.add_run()
            run.text = content
            run.font.size = Pt(font_size_pt)
            run.font.name = FONT_NAME_BODY
            run.font.color.rgb = COLOR_TEXT_MAIN
            
            if is_bold:
                run.font.bold = True
                run.font.color.rgb = COLOR_MAIN


def _draw_flow_horizontal(slide, x, y, w, h, label, text):
    # 1. Outer Container (Same style as normal section)
    # Border Box (Light Gray)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(x), Cm(y),
        Cm(w), Cm(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BORDER
    box.line.width = Pt(1.0)

    header_h = 1.0
    
    # Accent Bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(x + 0.1), Cm(y + 0.15),
        Cm(0.15), Cm(header_h - 0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_MAIN
    bar.line.fill.background()
    
    # Label
    label_box = slide.shapes.add_textbox(
        Cm(x + 0.4), Cm(y),
        Cm(w - 0.4), Cm(header_h)
    )
    p = label_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = FONT_NAME_BOLD
    p.font.color.rgb = COLOR_MAIN
    label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 2. Flow Content
    content_y = y + header_h + 0.2
    content_h = h - header_h - 0.4
    content_w = w - 0.8 # Padding
    start_x = x + 0.4
    
    # Extract steps from text (bullet points)
    # Assume lines starting with bullets are steps
    raw_lines = text.split('\n')
    steps = [line.strip().lstrip('・-●').strip() for line in raw_lines if line.strip()]
    
    if not steps: return

    # Calculate dimensions
    step_count = len(steps)
    # Max steps to display horizontally to avoid crowding: 4
    if step_count > 4: step_count = 4 
    display_steps = steps[:step_count]
    
    # Widths: [Box] [Arrow] [Box] ...
    # Arrow assume 0.8cm
    arrow_w = 0.8
    total_arrow_w = arrow_w * (step_count - 1)
    
    if total_arrow_w >= content_w:
        # Fallback to simple text if not enough space
        _draw_section(slide, x, y, w, h, label, text)
        return

    box_w = (content_w - total_arrow_w) / step_count
    
    curr_x = start_x
    for i, step_text in enumerate(display_steps):
        # Draw Box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(curr_x), Cm(content_y),
            Cm(box_w), Cm(content_h)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(240, 248, 255) # Light AliceBlue
        step_box.line.color.rgb = COLOR_MAIN
        step_box.line.width = Pt(1.0)
        
        # Text inside box
        tf = step_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_text
        p.font.size = Pt(10) # Smaller font for flow boxes
        p.font.name = FONT_NAME_BODY
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.alignment = PP_ALIGN.CENTER
        
        curr_x += box_w
        
        # Draw Arrow (if not last)
        if i < step_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Cm(curr_x + 0.1), Cm(content_y + content_h/2 - 0.2), # Center vertically
                Cm(arrow_w - 0.2), Cm(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.fill.background()
            
            curr_x += arrow_w
